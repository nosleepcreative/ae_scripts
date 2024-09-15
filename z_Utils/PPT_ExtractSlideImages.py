# https://www.freepdfconvert.com/ppt-to-pdf
# workflow for animating PPT
# mask each element of each page -> MasksToLayer.jsx
# Animate + Rift
# # Remember to install the python-pptx and Pillow libraries before running this script. 
#You can install them using pip: pip install python-pptx Pillow
from pptx import Presentation
myFile = r"C:\Users\Desmond\Downloads\_test\test.pptx"
pres = Presentation(myFile)
slide = pres.slides[0]
shape = slide.shapes[0]
image = shape.image
blob = image.blob
ext = image.ext
with open(f'image.{ext}', 'wb') as file:
    file.write(blob)